from pptx import Presentation
from pptx.util import Inches
from dateutil.relativedelta import relativedelta
import matplotlib
matplotlib.use('Agg')
import matplotlib.font_manager as fm
import sys
import os
import io
import pymysql
import pandas as pd
import numpy as np
import datetime

from dateutil import rrule
from matplotlib import pyplot as plt
from tensorflow.python.keras.models import Sequential, load_model
from tensorflow.python.keras.layers import Dense, LSTM, Dropout
from tensorflow.python.keras.callbacks import EarlyStopping, ModelCheckpoint
# tensorflow WARNING 제거
# tf.compat.v1.logging.set_verbosity(tf.compat.v1.logging.ERROR)


class Zabbix():
    def __init__(self, user, passwd, host, db, port, charset):
        self.db_connect = pymysql.connect(user=user, passwd=passwd, host=host, db=db, port=port, charset=charset)
        self.curs = self.db_connect.cursor()

    def read_db(self, sql_command):
        self.curs.execute(sql_command)
        df = pd.DataFrame(self.curs.fetchall())
        return df

    def read_one(self, sql_command):
        self.curs.execute(sql_command)
        one = self.curs.fetchone()
        return one

    def manipulate_db(self, sql_command):
        self.curs.execute(sql_command)
        self.db_connect.commit()

# def read_db(itemid):
#     """
#     :param: itemid (history_uint 테이블)
#     :return: df (columns: itemid, datetime, traffic)
#     """
#     db_connect = pymysql.connect(user='zabbix', passwd='zabbix', host='210.121.218.5', db='ZABBIXDB', port=3306,
#                                charset='euckr')
#     curs = db_connect.cursor(pymysql.cursors.DictCursor)
#     sql_command = "SELECT itemid, from_unixtime(clock) as datetime, value as traffic FROM history_uint WHERE itemid IN ({})".format(itemid)
#     curs.execute(sql_command)
#     df = pd.DataFrame(curs.fetchall())
#
#     return df


def read_db_thread(itemid, table_name):
    # """
    # :param: itemid (history_uint 테이블)
    # :return: df (columns: itemid, datetime, traffic)
    # """
    # global total_df
    # db_connect = pymysql.connect(user='zabbix', passwd='zabbix', host='210.121.218.5', db='ZABBIXDB', port=3306,
    #                              charset='euckr')
    # curs = db_connect.cursor(pymysql.cursors.DictCursor)
    # lock.acquire()
    # sql_command = "SELECT itemid, from_unixtime(clock) as datetime, value as traffic FROM {} WHERE itemid IN ({})".format(
    #     table_name, itemid)
    # curs.execute(sql_command)
    # df = pd.DataFrame(curs.fetchall())
    # # print(table_name)
    # # print(df)
    # total_df = pd.concat([total_df, df])
    # lock.release()
    # # print('-----total_df-----')
    # # print(total_df)
    print('itemid: {}'.format(itemid))
    print('table_name: {}'.format(table_name))


def insert_db(df):
    """
    예측값 zabbix db에 업로드
    """
    db_connect = pymysql.connect(user='zabbix', passwd='zabbix', host='210.121.218.5', db='ZABBIXDB', port=3306,
                                 charset='euckr')
    curs = db_connect.cursor(pymysql.cursors.DictCursor)
    for row in range(len(df)):
        sql_command = "INSERT INTO test_ys (itemid, dates, traffic) VALUES ('{}', '{}', '{}')".format(df.iloc[row][0],
                                                                                                        df.iloc[row][1],
                                                                                                        df.iloc[row][2])
        curs.execute(sql_command)
        db_connect.commit()


def delete_db(itemid):
    """
    zabbix db 업로드 전 해당 itemid row 삭제
    """
    db_connect = pymysql.connect(user='zabbix', passwd='zabbix', host='210.121.218.5', db='ZABBIXDB', port=3306,
                                 charset='euckr')
    curs = db_connect.cursor(pymysql.cursors.DictCursor)
    sql_command = "DELETE FROM test_ys WHERE itemid='{}'".format(itemid)
    curs.execute(sql_command)
    db_connect.commit()


def reframe_df(df):
    """
    :param db에서 로드한 데이터프레임
    :return: 재구성한 데이터프레임(index: date, column: 'dt2sin', 'traffic')
    """
    week2sec = 7 * 24 * 60 * 60   # 일주일을 초 단위로 변환
    dt2ts = df['datetime'].map(datetime.datetime.timestamp)   # datetime to timestamp
    df['dt2sin'] = np.sin(dt2ts*(2*np.pi/week2sec))   # datetime to sin
    df['date'] = df['datetime'].dt.date   # date column 추가
    df.drop(['itemid', 'datetime'], axis=1, inplace=True)
    df = df[['date', 'dt2sin', 'traffic']]
    dp_df = df.groupby(['date'], as_index=True).max()   # day peak dataframe
    dp_df = dp_df.dropna()
    return dp_df


class Scaler():
    def __init__(self, df, feature):
        self.df = df
        self.feature = feature
        self.max = max(self.df[self.feature])
        self.min = min(self.df[self.feature])

    def normalization(self):
        self.df[self.feature] = list((self.df[self.feature]-self.min)/(self.max-self.min))
        return self.df

    def rev_normalization(self, array):
        pred_array = array*(self.max-self.min)+self.min
        return pred_array


class Train():
    def __init__(self, itemid, in_steps, out_steps, valid_per, epochs, batch_size, unit, drop_per):
        self.itemid = itemid
        self.in_steps = in_steps
        self.out_steps = out_steps
        self.valid_per = valid_per
        self.epochs = epochs
        self.batch_size = batch_size
        self.unit = unit
        self.drop_per = drop_per

    def split_data(self, df):
        train_x = []
        train_y = []
        for i in range(len(df) - (self.in_steps + self.out_steps)):
            train_x.append(df.values[i: i + self.in_steps, :])
            train_y.append(df.values[i + self.in_steps: i + (self.in_steps + self.out_steps), -1])
        train_x, train_y = np.array(train_x), np.array(train_y)

        return train_x, train_y

    def generate_lstm(self, data):
        # data -> train_x
        model = Sequential([
            LSTM(self.unit, return_sequences=True, input_shape=(data.shape[1], data.shape[2])),
            Dropout(self.drop_per),
            LSTM(self.unit, return_sequences=False),
            Dropout(self.drop_per),
            Dense(self.out_steps)
        ])
        model.compile(optimizer='rmsprop', loss='mean_squared_error')
        return model

    def predict(self, df):
        train_x, train_y = self.split_data(df)
        model = self.generate_lstm(train_x)
        # verbose 1: 언제 training 멈추었는지 확인 가능
        early_stopping = EarlyStopping(monitor='loss', mode='min', patience=10, verbose=1)
        model_name = '{}.h5'.format(self.itemid)
        model_path = os.getcwd() + '/models/' + model_name
        model_check = ModelCheckpoint(filepath=model_path, monitor='loss', mode='min', save_best_only=True)
        # verbose 0: silence, 1: progress bar, 2: one line per each
        hist = model.fit(train_x, train_y, self.batch_size, self.epochs, validation_split=self.valid_per,
                         callbacks=[early_stopping, model_check], verbose=0)
        test_data = df.values[-self.in_steps:]
        test_x = test_data.reshape(1, self.in_steps, df.shape[1])
        # model = load_model('./models/{}.h5'.format(self.itemid))
        prediction = model.predict(test_x)
        return prediction


def making_report(datetimes, line_no, date_start, date_end):

    print('report 생성 시작...\n')

    line_num = line_no.replace('-', '')

    # Zabbix 클래스 개체 'zabbix' 생성
    user = 'zabbix'
    passwd = 'zabbix'
    host = '210.121.218.5'
    db = 'ZABBIXDB'
    port = 3306
    charset = 'euckr'
    zabbix = Zabbix(user=user, passwd=passwd, host=host, db=db, port=port, charset=charset)


    # neoss table 에서 데이터 가져오기
    sql_command = 'SELECT service_id, host_name, interface, ethernet_ip FROM neoss WHERE leased_line_num="{}"'.format(line_num)
    neoss_df = zabbix.read_db(sql_command)
    neoss_df.columns = ['service_id', 'host_name', 'interface', 'ethernet_ip']

    print('----- neoss_df -----')
    print(neoss_df, '\n')


    # neoss_df 에서 ip 데이터 뽑기. ppt #3  'ip 할당 현황' 데이터
    ethernet_ip = neoss_df['ethernet_ip'].values

    ip_count = [0 for _ in range(7)]   # ip count 하는 list [256, 128, 64, 32, 16, 8, 4] 순서
    total_ip = 0
    for ip in ethernet_ip:
        try:
            ip_end = int(ip[-2:])   # ip 마지막 두 자리 정수로 변경
            ip_count[ip_end - 24] += 1   # ip 마지막 두 자리 정수 변환. 24 -> 256, 25 -> 128, 24 -> 64 ...
            total_ip += pow(2, 32 - ip_end)
        except:
            pass
    ethernet_ip_24 = ip_count[0]
    ethernet_ip_25 = ip_count[1]
    ethernet_ip_26 = ip_count[2]
    ethernet_ip_27 = ip_count[3]
    ethernet_ip_28 = ip_count[4]
    ethernet_ip_29 = ip_count[5]
    ethernet_ip_30 = ip_count[6]

    ethernet_ip_list = ['-' for i in range(8)]
    for i in range(len(ethernet_ip)):
        if i < 8:
            start_ip = int(ethernet_ip[i].split('.')[-1].split('/')[0])
            ip_num = pow(2, 32-int(ethernet_ip[i][-2:]))
            ethernet_ip_list[i] = ethernet_ip[i][:-3] + '~' + str(start_ip + ip_num - 1)
        if i >= 8:
            break

    # neoss_df 에서 service_id 값 뽑기
    service_id = neoss_df['service_id'][0]

    # bidw table 에서 데이터 가져오기
    sql_command = 'SELECT industry_2, customer_name, contract_speed, service_name, contract_end FROM bidw WHERE ' \
                     'service_id="{}"'.format(service_id)
    bidw_df = zabbix.read_db(sql_command)
    bidw_df.columns = ['industry_2', 'customer_name', 'contract_speed', 'service_name', 'contract_end']
    print('----- bidw_df -----')
    print(bidw_df, '\n')

    # bidw table 에서 industry 기준 contract_speed
    industry = bidw_df['industry_2'][0]
    sql_command = 'SELECT industry_2, contract_speed FROM bidw WHERE industry_2="{}"'.format(industry)
    industry_speed_df = zabbix.read_db(sql_command)
    industry_speed_df.columns = ['industry_2', 'contract_speed']
    print('----- industry_speed_df -----')
    print(industry_speed_df, '\n')

    # contract_speed 값 변환 (M -> pow(10,6), G -> pow(10,9))
    speed_data = []

    for i in industry_speed_df['contract_speed']:
        if i[-1] == 'M':
            speed_data.append(int(i[:-1])*pow(10, 6))   # 수정
        elif i[-1] == 'G':
            speed_data.append(int(i[:-1])*pow(10, 9))   # 수정

    # speed 개수 count
    over_1G = 0
    over_500M = 0
    over_100M = 0
    under_100M = 0

    for i in speed_data:   # 수정
        if i >= pow(10, 9):
            over_1G += 1
        elif pow(10, 9) > i >= 5*pow(10, 8):
            over_500M += 1
        elif 5*pow(10, 8) > i >= pow(10, 8):
            over_100M += 1
        else:
            under_100M += 1

    print("1G 이상: {} 개".format(over_1G))
    print("500M 이상: {} 개".format(over_500M))
    print("100M 이상: {} 개".format(over_100M))
    print("100M 미만: {} 개".format(under_100M), '\n')


    # bidw table 과 neoss table 에서 service_id 기준 industry, customer_name, ethernet_ip 합치기. 동종업계 ip 비교 위함
    sql_command = 'SELECT industry_2, customer_name, service_id FROM bidw WHERE industry_2="{}"'\
        .format(industry)
    merge_1 = zabbix.read_db(sql_command)
    merge_1.columns = ['industry_2', 'customer_name', 'service_id']
    print('----- merge_1 -----')
    print(merge_1, '\n')

    sql_command = 'SELECT service_id, ethernet_ip FROM neoss'
    merge_2 = zabbix.read_db(sql_command)
    merge_2.columns = ['service_id', 'ethernet_ip']
    print('----- merge_2 -----')
    print(merge_2, '\n')

    merge_df = pd.merge(merge_1, merge_2, on="service_id")
    print('----- merge_df -----')
    print(merge_df, '\n')

    # 동종업계 고객 별 ip value 설정
    customer_name_dict = {}
    for customer in merge_df['customer_name']:
        customer_name_dict[customer] = 0

    for row in merge_df.values:   # col1: industry, col2: customer_name, col3: service_id, col4: ethernet_ip
        try:
            customer_name_dict[row[1]] += pow(2, 32 - int(row[3][-2:]))
        except:
            pass

    # 동종업계 ip 별 count
    ip_list = [1000, 500, 256, 128, 64, 32, 16, 8, 4]
    customer_ip_count = [0 for _ in range(9)]

    for ip in customer_name_dict.values():
        for index, value in enumerate(ip_list):
            if ip >= ip_list[index]:
                customer_ip_count[index] += 1
                break
    ip_1000 = customer_ip_count[0]
    ip_500 = customer_ip_count[1]
    ip_256 = customer_ip_count[2]
    ip_128 = customer_ip_count[3]
    ip_64 = customer_ip_count[4]
    ip_32 = customer_ip_count[5]
    ip_16 = customer_ip_count[6]
    ip_8 = customer_ip_count[7]
    ip_4 = customer_ip_count[8]

    print('----- customer_ip_count -----')
    print(customer_ip_count, '\n')


    # host_name 과 hostid 추출

    host_name = neoss_df['host_name'][0]
    print('host_name: {}'.format(host_name))

    sql_command = 'SELECT hostid FROM hosts WHERE host="{}"'.format(host_name)
    hostid = zabbix.read_one(sql_command)
    hostid = hostid[0]
    print('hostid: {}'.format(hostid), '\n')


    # itemid 추출

    # interface = '%' + neoss_df['interface'][0] + '%'   # 수정: df_neoss -> neoss_df
    # 정규식 표현. 괄호 안의 문자열 포함 및 뒤에 숫자 제외.
    interface = neoss_df['interface'][0]
    in_interface = '(\\In.+' + interface + ')+[^0-9]'
    out_interface = '(\\Out.+' + interface + ')+[^0-9]'
    interface_list = [in_interface, out_interface]
    print('----- interface_list -----')
    print('interface: {}'.format(interface))
    print(interface_list, '\n')

    itemid_list = []
    for interface in interface_list:
        sql_command = 'SELECT itemid FROM items WHERE hostid="{}" and key_ REGEXP "{}"'.format(str(hostid), interface)
        itemid = zabbix.read_one(sql_command)   # tuple 형식의 itemid
        print('itemid: {}'.format(itemid[0]), '\n')
        if itemid == None:
            print('itemid를 찾을 수 없습니다.')
            sys.exit()
        else:
            itemid_list.append(itemid[0])   # [0]: in_itemid, [1]: out_itemid

    print('----- itemid list -----')
    print(itemid_list, '\n')


    # 제공속도와 청약속도 추출
    sql_command = 'SELECT engre, gbic FROM realspeed WHERE leased_line_num={}'.format(line_num)
    df_system = zabbix.read_db(sql_command)
    df_system.columns = ['engre', 'gbic']
    print('----- df_system -----')
    print(df_system, '\n')

    # offer_speed: 제공속도, contract_speed: 청약속도, _M: 각 속도 메가 단위
    offer_speed = int(df_system['engre'][0])
    offer_speed_M = int(offer_speed/pow(10, 6))
    contract_speed = int(df_system['gbic'][0])
    contract_speed_M = int(contract_speed/pow(10, 6))

    if int(offer_speed) >= contract_speed:
        if offer_speed >= 1000*pow(10, 6):
            script2_29 = '시설 변경 필요'
        else:
            script2_29 = '증속 불가'
        script2_30 = '중단'
    else:
        script2_29 = '시설변경 없이\n' \
                     '추가 ' + str(contract_speed_M - offer_speed_M) + 'M 증속 가능'
        script2_30 = '무중단 작업 가능'


    # AI 학습 및 예측
    for itemid in itemid_list:
        sql_command = "SELECT itemid, from_unixtime(clock) as datetime, value as traffic FROM history_uint WHERE itemid IN ({})".format(itemid)
        df = zabbix.read_db(sql_command)
        df.columns = ['itemid', 'datetime', 'traffic']
        print(df)
        dp_df = reframe_df(df)
        traffic_scaler = Scaler(dp_df, 'traffic')
        sc_df = traffic_scaler.normalization()
        df_length = len(sc_df)
        print('df_length: {}'.format(df_length))
        in_steps = (df_length // 42) * 7
        print('in_steps: {}'.format(in_steps))
        out_steps = in_steps
        print('out_steps: {}'.format(out_steps))
        if in_steps == 0:
            out_steps = 7

        # 해당 itemid에 해당하는 df 생성
        upload_df = pd.DataFrame(index=range(0, out_steps), columns=['itemid', 'dates', 'traffic'])
        upload_df['itemid'] = itemid
        # 예측하고자 하는 첫날부터 out_steps 만큼의 list 생성
        upload_df['dates'] = [(sc_df.index[-1] + datetime.timedelta(days=i)).strftime('%y%m%d') for i in range(out_steps)]
        if in_steps != 0:
            trainer = Train(itemid, in_steps, out_steps, valid_per=0, epochs=100, batch_size=128, unit=128, drop_per=0.1)
            # 음의 예측값은 0으로 변환
            real_prediction = traffic_scaler.rev_normalization(trainer.predict(sc_df)[0])
            upload_df['traffic'] = [max(num, 0) for num in real_prediction]
            print('예측 완료!', '\n')
        else:
            print('데이터 부족으로 예측 불가', '\n')
            upload_df['traffic'] = 0
        # upload_df.to_csv('C:/Users/User/Desktop/upload_df.csv', index=False)
        delete_db(itemid)
        insert_db(upload_df)


    ###### 스레드 테스트 ######
    # itemid = " "
    # db에서 읽어온 date값의 앞 7자리 ("YYYY-MM")
    date_starts = date_start[:7]
    date_ends = date_end[:7]

    # string을 datetime으로 바꿔주기
    start_date = datetime.datetime.strptime(date_starts, '%Y-%m')
    end_date = datetime.datetime.strptime(date_ends, '%Y-%m')

    # 월별 datetime list 생성
    months_list = list(rrule.rrule(rrule.MONTHLY, dtstart=start_date, until=end_date))

    # 'YYYY-MM' 꼴의 시작 날짜, 마지막 날짜 생성
    start = str(months_list[0].year) + '_' + str(months_list[0].month).zfill(2)
    end = str(months_list[-1].year) + '_' + str(months_list[-1].month).zfill(2)

    # 시작 날짜와 마지막 날짜 사이의 월별 날짜 생성
    date_list = [(str(obj.year) + '_' + str(obj.month).zfill(2)) for obj in months_list]

    # 월별 날짜를 포함한 table명 list 생성
    table = 'history_uint_'
    table_list = [(table + date) for date in date_list]

    # lock = threading.Lock()
    # total_df = pd.DataFrame(columns=['itemid', 'datetime', 'traffic'])
    # for table in table_list:
    #     my_thread = threading.Thread(target=read_db_thread, args=(itemid, table,))
    #     my_thread.start()


    # 실제 traffic 값 추출
    real_traffic_list = []
    for itemid in itemid_list:
        sql_command = 'SELECT itemid, from_unixtime(clock), value  FROM history_uint WHERE itemid="{}" ' \
                         'and "{}" <= from_unixtime(clock) AND from_unixtime(clock) <= "{}"'\
                        .format(itemid, date_start, date_end)
        traffic_df = zabbix.read_db(sql_command)
        traffic_df.columns = ['itemid', 'clock', 'traffic']
        real_traffic_list.append(traffic_df)
    in_real_traffic_df = real_traffic_list[0]
    out_real_traffic_df = real_traffic_list[1]
    print('-----in_real_traffic_df-----')
    print(in_real_traffic_df, '\n')
    print('-----out_real_traffic_df-----')
    print(out_real_traffic_df, '\n')


    # 트래픽 임계치 count 행렬
    # 1행: in, 2행: out
    # 0열~6열: 제공속도 100% 초과/ 90~100% / 80~90% / 70~80% / 50~70% / 49% 이하
    over_count = np.zeros((2, 6), dtype='int')

    for index, df in enumerate(real_traffic_list):
        for value in df['traffic']:
            try:
                if int(value) > int(offer_speed):
                    over_count[index, 0] += 1
                elif int(offer_speed) > int(value) >= int(offer_speed)*0.9:
                    over_count[index, 1] += 1
                elif int(offer_speed)*0.9 > int(value) >= int(offer_speed)*0.8:
                    over_count[index, 2] += 1
                elif int(offer_speed)*0.8 > int(value) >= int(offer_speed)*0.7:
                    over_count[index, 3] += 1
                elif int(offer_speed)*0.7 > int(value) >= int(offer_speed)*0.5:
                    over_count[index, 4] += 1
                else:
                    over_count[index, 5] += 1
            except:
                pass

    in_over_100 = over_count[0, 0]
    in_over_90 = over_count[0, 1]
    in_over_80 = over_count[0, 2]
    in_over_70 = over_count[0, 3]
    in_over_50 = over_count[0, 4]
    in_over_49 = over_count[0, 5]
    out_over_100 = over_count[1, 0]
    out_over_90 = over_count[1, 1]
    out_over_80 = over_count[1, 2]
    out_over_70 = over_count[1, 3]
    out_over_50 = over_count[1, 4]
    out_over_49 = over_count[1, 5]

    if in_over_100 > 0:
        script1_12_0 = '업/다운 트래픽이 제공속도인 ' + str(offer_speed_M) + 'Mbps를 초과하는 트래픽 폭주 ' + \
                       str(in_over_100) + '회 발생'
        script1_12_1 = '5분 평균치 감안 시 실제 체감속도는 더욱 늦을 것으로 판단됨'
        script1_12_2 = '정확한 원인은 고객사 보안장비 등에서 IP/서비스별 확인 필요'
        script1_12 = [script1_12_0, script1_12_1, script1_12_2]

        script1_14_0 = '정상적인 업무 트래픽에 의한 폭주일 경우, GiGA Office 도입 또는 코넷 증속 필요'
        script1_14_1 = '안정적인 인터넷 이용을 위해 QoS 기반 UTM/방화벽 망구성 및 트래픽 제한 등 관리 필요'
        script1_14_2 = '네트워크 및 보안 진단 컨설팅 필요'
        script1_14 = [script1_14_0, script1_14_1, script1_14_2]
    elif in_over_90 > 0:
        script1_12_0 = '업/다운 트래픽이 제공속도 ' + str(offer_speed_M) + 'Mbps에 육박하는 트래픽 폭주 ' +\
                     str(in_over_90) + '회 발생'
        script1_12_1 = '5분 평균치 감안 시 실제 체감속도는 더욱 늦을 것으로 판단됨'
        script1_12_2 = '정확한 원인은 고객사 보안장비 등에서 IP/서비스별 확인 필요'
        script1_12 = [script1_12_0, script1_12_1, script1_12_2]

        script1_14_0 = '정상적인 업무 트래픽에 의한 폭주일 경우, GiGA Office 도입 또는 코넷 증속 권고'
        script1_14_1 = '안정적인 인터넷 이용을 위해 QoS 기반 UTM/방화벽 망구성 및 트래픽 제한 등 관리 필요'
        script1_14_2 = '네트워크 및 보안 진단 컨설팅 필요'
        script1_14 = [script1_14_0, script1_14_1, script1_14_2]
    elif (in_over_70 + in_over_80) > 0:
        script1_12_0 = '다운로드의 경우 권고치를 초과하는 트래픽 폭주 ' + str(in_over_70 + in_over_80) + '회 발생'
        script1_12_1 = '5분 평균치 감안 시 실제 체감속도는 다소 늦을 수 있음'
        script1_12 = [script1_12_0, script1_12_1]

        script1_14_0 = '트래픽 증가를 고려하여 GiGA Office 도입 또는 코넷 증속 권고'
        script1_14_1 = '주기적인 모니터링으로 트래픽 추적 관리 권고'
        script1_14 = [script1_14_0, script1_14_1]
    elif in_over_50 > 0:
        script1_12_0 = '업/다운 트래픽이 권고치 이내로서 비교적 안정적임'
        script1_12_1 = '5분 평균 집계치임을 감안하더라도 체감속도는 비교적 안정적일 것으로 판단됨'
        script1_12 = [script1_12_0, script1_12_1]

        script1_14_0 = '주기적인 모니터링으로 트래픽 추적 관리 권고'
        script1_14 = [script1_14_0]
    elif in_over_49 > 0:
        script1_12_0 = '업/다운 트래픽이 권고치 이내로 매우 안정적임'
        script1_12 = [script1_12_0]

        script1_14_0 = '주기적인 모니터링으로 트래픽 추적 관리 권고'
        script1_14 = [script1_14_0]
    else:
        script1_12_0 = '트래픽 이용량이 거의 없음'
        script1_12 = [script1_12_0]

        script1_14_0 = '유휴 또는 백업용 회선으로 추정됨'
        script1_14 = [script1_14_0]


    # AI 예측 데이터 추출
    pred_traffic_list = []
    for itemid in itemid_list:
        sql_command = 'SELECT itemid, dates, traffic FROM test_ys WHERE itemid="{}"'.format(itemid)
        pred_traffic_df = zabbix.read_db(sql_command)
        pred_traffic_df.columns = ['itemid', 'dates', 'traffic']
        pred_traffic_list.append(pred_traffic_df)
    in_pred_traffic_df = pred_traffic_list[0]
    out_pred_traffic_df = pred_traffic_list[1]
    print('----- in_pred_traffic_df -----')
    print(in_pred_traffic_df, '\n')
    print('----- out_pred_traffic_df -----')
    print(out_pred_traffic_df, '\n')

    dates = ['20' + i[:2] + '-' + i[2:4] + '-' + i[4:] for i in pred_traffic_df['dates']]

    # 예측치와 제공속도 비교
    pred_over_count = np.zeros((2, 3), dtype='int')

    for index, df in enumerate(pred_traffic_list):
        for traffic in df['traffic']:
            try:
                if int(traffic) > int(offer_speed)*0.9:
                    pred_over_count[index, 0] += 1
                elif int(offer_speed)*0.9 > int(traffic) >= int(offer_speed)*0.7:
                    pred_over_count[index, 1] += 1
                else:
                    pred_over_count[index, 2] += 1
            except:
                pass

    pred_in_over_90 = pred_over_count[0, 0]
    pred_in_over_70 = pred_over_count[0, 1]
    pred_in_under_70 = pred_over_count[0, 2]
    pred_out_over_90 = pred_over_count[1, 0]
    pred_out_over_70 = pred_over_count[1, 1]
    pred_out_under_70 = pred_over_count[1, 2]

    if (pred_in_over_90 + pred_out_over_90) > 0:
        script1_13_0 = 'AI 예측 알고리즘을 통한 분석결과 지속적인 폭주가 예상 됨'
        script1_13_1 = '향후 권고치' + str(int(offer_speed_M*0.7)) + 'M를 초과 할 것으로 예측 됨'
        script1_13 = [script1_13_0, script1_13_1]
    elif (pred_in_over_70 + pred_out_over_70) > 0:
        script1_13_0 = 'AI 예측 알고리즘을 통한분석결과 주기적 폭주가 예상 됨'
        script1_13 = [script1_13_0]
    else:
        script1_13_0 = 'AI 예측 알고리즘을 통한 분석결과 완만한 증가가 예상 됨'
        script1_13 = [script1_13_0]


    #########total script###########

    # -------------------한글 폰트 지원-------------------
    path_gothic = 'C:/Windows/Fonts/malgunbd.ttf'
    fontprop1 = fm.FontProperties(fname=path_gothic, size=10)

    # -----------------------페이지별 항목----------------------------
    prs_test = Presentation('Sales_More_ppt2.2.pptx')
    slide0 = prs_test.slides.add_slide(prs_test.slide_layouts[0])
    # for shape in slide0.placeholders:
    #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
    # slide1 = prs_test.slides.add_slide(prs_test.slide_layouts[1])
    # for shape in slide1.placeholders:
    #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
    # slide2 = prs_test.slides.add_slide(prs_test.slide_layouts[2])
    # for shape in slide2.placeholders:
    #     print('%d %s' % (shape.placeholder_format.idx, shape.name))


    data0_1 = bidw_df['customer_name'][0]
    data0_2 = str(datetime.datetime.now().date())

    data1_1 = bidw_df['customer_name'][0]
    data1_2 = neoss_df['service_id'][0]   # 수정: df_neoss -> neoss_df
    data1_3 = line_num[:2] + '-' + line_num[2:8] + '-' + line_num[-4:]
    data1_4 = str(offer_speed_M) + 'M'
    data1_5 = date_start
    data1_6 = date_end
    data1_7 = dates[0]
    data1_8 = dates[-1]
    data1_9 = str(int(in_real_traffic_df['traffic'].max()/pow(10, 6))) + 'M / ' + str(int(out_real_traffic_df['traffic'].max()/pow(10, 6))) + 'M'
    data1_10 = format((in_real_traffic_df['traffic'].max()/offer_speed)*100, '.1f') + '% / ' + format((out_real_traffic_df['traffic'].max()/offer_speed)*100, '.1f') + '%'
    data1_11 = str(int(in_pred_traffic_df['traffic'].max()/pow(10, 6))) + 'M / ' + str(int(out_pred_traffic_df['traffic'].max()/pow(10, 6))) + 'M'
    data1_12 = format((in_pred_traffic_df['traffic'].max()/offer_speed)*100, '.1f') + '% / ' + format((out_pred_traffic_df['traffic'].max()/offer_speed)*100, '.1f') + '%'

    # 약정만료 기한 표시
    contract_end = bidw_df['contract_end'][0]
    contract_ends = datetime.datetime.strptime(contract_end, '%Y%m') + relativedelta(months=1) - datetime.timedelta(days=1)
    now = datetime.datetime.strptime(datetimes, '%Y-%m-%d-%H%M%S')

    data2_1 = bidw_df['service_name'][0]
    # data2_2 = bidw_df['contract_speed'][0]   청약속도
    data2_3 = str(offer_speed_M) + 'M'
    if contract_ends > now:
        data2_4 = contract_end[:4] + '년 ' +  contract_end[4:] + '월'
    else:
        data2_4 = '-'
    data2_5 = str(ethernet_ip_24)
    data2_6 = str(ethernet_ip_25)
    data2_7 = str(ethernet_ip_26)
    data2_8 = str(ethernet_ip_27)
    data2_9 = str(ethernet_ip_28)
    data2_10 = str(ethernet_ip_29)
    data2_11 = ethernet_ip_list[0]
    data2_12 = ethernet_ip_list[1]
    data2_13 = ethernet_ip_list[2]
    data2_14 = ethernet_ip_list[3]
    data2_15 = ethernet_ip_list[4]
    data2_16 = ethernet_ip_list[5]
    data2_17 = ethernet_ip_list[6]
    data2_18 = ethernet_ip_list[7]

    data2_19 = str(in_over_70 + out_over_70)
    data2_20 = str(in_over_80 + out_over_80)
    data2_21 = str(in_over_90 + in_over_100 + out_over_90 + out_over_100)
    data2_22 = str(int(data2_19) + int(data2_20) + int(data2_21))
    data2_23 = '2_23'
    data2_24 = '2_24'
    data2_25 = '2_25'
    data2_26 = '2_26'
    data2_27 = str(offer_speed_M) + 'M'
    data2_28 = str(contract_speed_M) + 'M'
    data2_29 = script2_29
    data2_30 = script2_30
    data2_31 = industry
    data2_32 = str(ethernet_ip_30)
    data2_33 = '2_33'
    data2_34 = '2_34'
    data2_35 = '2_35'
    data2_36 = '2_36'
    data2_37 = '2_37'
    data2_38 = '2_38'
    data2_39 = '2_39'
    data2_40 = '2_40'

    # ppt 파일 불러오기
    prs = Presentation('Sales_More_ppt2.2.pptx')
    # 슬라이드 마스터 첫번째 레이아웃 불러오기 (prs.slide_layouts)
    slide_layout = prs.slide_layouts[0]
    # 슬라이드 생성하기
    slide = prs.slides.add_slide(slide_layout)
    # 슬라이드 마스터의 개체 틀
    shapes = slide.shapes

    #####page0####
    page0_1_frame = shapes.placeholders[20].text_frame
    page0_1_frame.text = data0_1
    page0_2_frame = shapes.placeholders[19].text_frame
    page0_2_frame.text = data0_2

    #####page1####
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes

    page1_1_frame = shapes.placeholders[14].text_frame
    page1_1_frame.text = data1_1
    page1_2_frame = shapes.placeholders[15].text_frame
    page1_2_frame.text = data1_2
    page1_3_frame = shapes.placeholders[16].text_frame
    page1_3_frame.text = data1_3
    page1_4_frame = shapes.placeholders[17].text_frame
    page1_4_frame.text = data1_4
    page1_5_frame = shapes.placeholders[29].text_frame
    page1_5_frame.text = data1_5
    page1_6_frame = shapes.placeholders[30].text_frame
    page1_6_frame.text = data1_6
    page1_7_frame = shapes.placeholders[31].text_frame
    page1_7_frame.text = data1_7
    page1_8_frame = shapes.placeholders[32].text_frame
    page1_8_frame.text = data1_8
    page1_9_frame = shapes.placeholders[18].text_frame
    page1_9_frame.text = data1_9
    page1_10_frame = shapes.placeholders[27].text_frame
    page1_10_frame.text = data1_10
    page1_11_frame = shapes.placeholders[33].text_frame
    page1_11_frame.text = data1_11
    page1_12_frame = shapes.placeholders[34].text_frame
    page1_12_frame.text = data1_12

    if len(script1_12) == 3:
        page1_12_frame = shapes.placeholders[22].text_frame.paragraphs[0]
        page1_12_frame.text = script1_12[0]
        page1_12_frame_1 = shapes.placeholders[22].text_frame.add_paragraph()
        page1_12_frame_1.text = script1_12[1]
        page1_12_frame_1.level = 1
        page1_12_frame_2 = shapes.placeholders[22].text_frame.add_paragraph()
        page1_12_frame_2.text = script1_12[2]
    elif len(script1_12) == 2:
        page1_12_frame = shapes.placeholders[22].text_frame.paragraphs[0]
        page1_12_frame.text = script1_12[0]
        page1_12_frame_1 = shapes.placeholders[22].text_frame.add_paragraph()
        page1_12_frame_1.text = script1_12[1]
        page1_12_frame_1.level = 1
    else:
        page1_12_frame = shapes.placeholders[22].text_frame.paragraphs[0]
        page1_12_frame.text = script1_12[0]

    if len(script1_13) == 2:
        page1_13_frame = shapes.placeholders[25].text_frame.paragraphs[0]
        page1_13_frame.text = script1_12[0]
        page1_13_frame_1 = shapes.placeholders[25].text_frame.add_paragraph()
        page1_13_frame_1.text = script1_12[1]
    else:
        page1_13_frame = shapes.placeholders[25].text_frame.paragraphs[0]
        page1_13_frame.text = script1_12[0]

    if len(script1_14) == 3:
        page1_14_frame = shapes.placeholders[26].text_frame.paragraphs[0]
        page1_14_frame.text = script1_14[0]
        page1_14_frame_1 = shapes.placeholders[26].text_frame.add_paragraph()
        page1_14_frame_1.text = script1_14[1]
        page1_14_frame_2 = shapes.placeholders[26].text_frame.add_paragraph()
        page1_14_frame_2.text = script1_14[2]
    elif len(script1_14) == 2:
        page1_14_frame = shapes.placeholders[26].text_frame.paragraphs[0]
        page1_14_frame.text = script1_14[0]
        page1_14_frame_1 = shapes.placeholders[26].text_frame.add_paragraph()
        page1_14_frame_1.text = script1_14[1]
    else:
        page1_14_frame = shapes.placeholders[26].text_frame.paragraphs[0]
        page1_14_frame.text = script1_14[0]


    ##### Traffic plot #####

    fig = (14, 2)
    plt.figure(figsize=fig)

    plt.title('M', loc="left", fontsize="10")
    plt.plot(in_real_traffic_df['traffic']/pow(10, 6), linestyle='-.', linewidth=1, color='blue', label='Up')
    plt.plot(out_real_traffic_df['traffic']/pow(10, 6), linestyle=':', linewidth=1, color='red', label='Down')
    plt.legend(loc='upper right')

    xticks_value = [i for i in range(0, len(in_real_traffic_df), (len(in_real_traffic_df) // 10))]
    xticks = [in_real_traffic_df['clock'][i].strftime('%Y%m%d%H%M')[2:8] for i in xticks_value]
    plt.xticks(xticks_value, xticks, rotation=45, fontsize=12)

    y_lim = max(max(in_real_traffic_df['traffic']), max(out_real_traffic_df['traffic'])) * 1.2
    plt.ylim(0, y_lim / pow(10, 6))

    if int(offer_speed) * 0.7 <= max(out_real_traffic_df['traffic']) * 1.2 < int(offer_speed):
        plt.axhline(y=int(offer_speed_M * 0.7), color='orange', linewidth=1, linestyle='--')
        plt.text(0, int(offer_speed_M * 0.7), "권고속도 : " + str(int(offer_speed_M * 0.7)) + 'M', fontsize=10, fontproperties=fontprop1)
        # plt.text(0, int(offer_speed_M * 0.7), "권고속도 : " + str(int(offer_speed_M * 0.7)) + 'M', fontsize=10)
    elif int(offer_speed) <= max(out_real_traffic_df['traffic']) * 1.2:
        plt.axhline(y=offer_speed_M, color='red', linewidth=1, linestyle='--')
        plt.text(0, offer_speed_M, "제공속도 : " + str(offer_speed_M) + 'M', fontsize=10, fontproperties=fontprop1)
        # plt.text(0, offer_speed_M, "제공속도 : " + str(offer_speed_M) + 'M', fontsize=10)
        plt.axhline(y=int(offer_speed_M * 0.7), color='orange', linewidth=1, linestyle='--')
        plt.text(0, int(offer_speed_M * 0.7), "권고속도 : " + str(int(offer_speed_M * 0.7)) + 'M', fontsize=10, fontproperties=fontprop1)
        # plt.text(0, int(offer_speed_M * 0.7), "권고속도 : " + str(int(offer_speed_M * 0.7)) + 'M', fontsize=10)

    image_stream_types = io.BytesIO()
    plt.savefig(image_stream_types, bbox_inches='tight', transparent=True)
    left = Inches(0.35)
    top = Inches(2.3)
    width = Inches(6)
    height = Inches(1.9)
    pic = shapes.add_picture(image_stream_types, left, top, width, height)
    plt.close()


    ##### AI plot #####

    if in_steps == 0:   # 데이터 부족 시 그림 삽입
        data_img_path = './images/data_lack.jpg'
        left_cm = 17.93
        left = Inches(left_cm / 2.54)  # inch = cm/2.54
        top_cm = 6.13
        top = Inches(top_cm / 2.54)
        pic = slide.shapes.add_picture(data_img_path, left, top)
    else:
        fig = (8, 2)
        plt.figure(figsize=fig)

        plt.title('M', loc="left", fontsize="10")
        plt.plot(in_pred_traffic_df['traffic']/pow(10, 6), linestyle='-.', linewidth=1, color='blue', label='Up')
        plt.plot(out_pred_traffic_df['traffic']/pow(10, 6), linestyle=':', linewidth=1, color='red', label='Down')
        plt.legend(loc='upper right')

        xticks_value = [i for i in range(0, len(in_pred_traffic_df), 1)]
        xticks = [in_pred_traffic_df['dates'][i] for i in xticks_value]
        plt.xticks(xticks_value, xticks, rotation=45, fontsize=12)

        plt.ylim(0, y_lim / pow(10, 6))   # 예측 그래프의 y축 limit은 실제 그래프와 동일

        if int(offer_speed) * 0.7 <= max(out_real_traffic_df['traffic']) * 1.2 < int(offer_speed):
            plt.axhline(y=int(offer_speed_M * 0.7), color='orange', linewidth=1, linestyle='--')
            plt.text(0, int(offer_speed_M * 0.7), "권고속도 : " + str(int(offer_speed_M * 0.7)) + 'M', fontsize=10, fontproperties=fontprop1)
            # plt.text(0, int(offer_speed_M * 0.7), "권고속도 : " + str(int(offer_speed_M * 0.7)) + 'M', fontsize=10)
        elif int(offer_speed) <= max(out_real_traffic_df['traffic']) * 1.2:
            plt.axhline(y=offer_speed_M, color='red', linewidth=1, linestyle='--')
            plt.text(0, offer_speed_M, "제공속도 : " + str(offer_speed_M) + 'M', fontsize=10, fontproperties=fontprop1)
            # plt.text(0, offer_speed_M, "제공속도 : " + str(offer_speed_M) + 'M', fontsize=10)
            plt.axhline(y=int(offer_speed_M * 0.7), color='orange', linewidth=1, linestyle='--')
            plt.text(0, int(offer_speed_M * 0.7), "권고속도 : " + str(int(offer_speed_M * 0.7)) + 'M', fontsize=10, fontproperties=fontprop1)
            # plt.text(0, int(offer_speed_M * 0.7), "권고속도 : " + str(int(offer_speed_M * 0.7)) + 'M', fontsize=10)

        image_stream_types = io.BytesIO()
        plt.savefig(image_stream_types, bbox_inches='tight', transparent=True)
        left = Inches(6.9)
        top = Inches(2.3)
        width = Inches(3.5)
        height = Inches(1.9)
        pic = shapes.add_picture(image_stream_types, left, top, width, height)
        plt.close()


    #####page2#####
    slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes

    page2_1_frame = shapes.placeholders[29].text_frame
    page2_1_frame.text = data2_1
    # page2_2_frame = shapes.placeholders[33].text_frame
    # page2_2_frame.text = data2_2
    page2_3_frame = shapes.placeholders[34].text_frame
    page2_3_frame.text = data2_3
    page2_4_frame = shapes.placeholders[39].text_frame
    page2_4_frame.text = data2_4
    page2_5_frame = shapes.placeholders[40].text_frame
    page2_5_frame.text = data2_5
    page2_6_frame = shapes.placeholders[41].text_frame
    page2_6_frame.text = data2_6
    page2_7_frame = shapes.placeholders[42].text_frame
    page2_7_frame.text = data2_7
    page2_8_frame = shapes.placeholders[43].text_frame
    page2_8_frame.text = data2_8
    page2_9_frame = shapes.placeholders[44].text_frame
    page2_9_frame.text = data2_9
    page2_10_frame = shapes.placeholders[45].text_frame
    page2_10_frame.text = data2_10
    page2_11_frame = shapes.placeholders[73].text_frame
    page2_11_frame.text = data2_11
    page2_12_frame = shapes.placeholders[74].text_frame
    page2_12_frame.text = data2_12
    page2_13_frame = shapes.placeholders[75].text_frame
    page2_13_frame.text = data2_13
    page2_14_frame = shapes.placeholders[76].text_frame
    page2_14_frame.text = data2_14
    page2_15_frame = shapes.placeholders[35].text_frame
    page2_15_frame.text = data2_15
    page2_16_frame = shapes.placeholders[36].text_frame
    page2_16_frame.text = data2_16
    page2_17_frame = shapes.placeholders[37].text_frame
    page2_17_frame.text = data2_17
    page2_18_frame = shapes.placeholders[72].text_frame
    page2_18_frame.text = data2_18
    page2_19_frame = shapes.placeholders[52].text_frame
    page2_19_frame.text = data2_19
    page2_20_frame = shapes.placeholders[53].text_frame
    page2_20_frame.text = data2_20
    page2_21_frame = shapes.placeholders[54].text_frame
    page2_21_frame.text = data2_21
    page2_22_frame = shapes.placeholders[67].text_frame
    page2_22_frame.text = data2_22
    page2_23_frame = shapes.placeholders[68].text_frame
    page2_23_frame.text = data2_23
    page2_24_frame = shapes.placeholders[69].text_frame
    page2_24_frame.text = data2_24
    page2_25_frame = shapes.placeholders[70].text_frame
    page2_25_frame.text = data2_25
    page2_26_frame = shapes.placeholders[71].text_frame
    page2_26_frame.text = data2_26
    page2_27_frame = shapes.placeholders[47].text_frame
    page2_27_frame.text = data2_27
    page2_28_frame = shapes.placeholders[48].text_frame
    page2_28_frame.text = data2_28
    page2_29_frame = shapes.placeholders[27].text_frame
    page2_29_frame.text = data2_29
    page2_30_frame = shapes.placeholders[28].text_frame
    page2_30_frame.text = data2_30
    page2_31_frame = shapes.placeholders[66].text_frame
    page2_31_frame.text = data2_31
    page2_32_frame = shapes.placeholders[77].text_frame
    page2_32_frame.text = data2_32
    page2_33_frame = shapes.placeholders[78].text_frame
    page2_33_frame.text = data2_33
    page2_34_frame = shapes.placeholders[79].text_frame
    page2_34_frame.text = data2_34
    page2_35_frame = shapes.placeholders[80].text_frame
    page2_35_frame.text = data2_35
    page2_36_frame = shapes.placeholders[81].text_frame
    page2_36_frame.text = data2_36
    page2_37_frame = shapes.placeholders[82].text_frame
    page2_37_frame.text = data2_37
    page2_38_frame = shapes.placeholders[83].text_frame
    page2_38_frame.text = data2_38
    page2_39_frame = shapes.placeholders[84].text_frame
    page2_39_frame.text = data2_39
    page2_40_frame = shapes.placeholders[85].text_frame
    page2_40_frame.text = data2_40


    ##### 서비스 중단 여부 동그라미 #####
    circle_img_path = './images/circle.png'
    if script2_30 == '무중단 작업 가능':
        left_cm = 7.86
    elif script2_30 == '중단':
        left_cm = 9.92
    left = Inches(left_cm/2.54)   # inch = cm/2.54
    top_cm = 13.29
    top = Inches(top_cm/2.54)
    pic = slide.shapes.add_picture(circle_img_path, left, top)


    fig = (2, 1)
    plt.figure(figsize=fig)
    plt.barh(1, contract_speed_M, color='#00C0AA')
    plt.axvline(x=contract_speed_M, color='red', linewidth=0.5, linestyle='--')
    plt.barh(1, offer_speed_M, color='#4C4C4E', alpha=0.8)
    plt.axvline(x=offer_speed_M, color='red', linewidth=0.5, linestyle='--')
    plt.axis('off')
    plt.ylim(0, 2)
    image_stream_types = io.BytesIO()
    plt.savefig(image_stream_types, bbox_inches='tight', transparent=True)
    left = Inches(0.8)
    top = Inches(5.2)
    width = Inches(1.6)
    height = Inches(0.8)
    pic = shapes.add_picture(image_stream_types, left, top, width, height)
    plt.close()

    #####동종업계 속도 이용현황#####

    if bidw_df['contract_speed'][0][-1] == 'G':
        contract_sp = int(bidw_df['contract_speed'][0][:-1]) * pow(10, 9)
    elif bidw_df['contract_speed'][0][-1] == 'M':
        contract_sp = int(bidw_df['contract_speed'][0][:-1]) * pow(10, 6)

    if contract_sp >= 1000000000:
        explode = [0.1, 0, 0, 0]
    elif 1000000000 > contract_sp >= 500000000:
        explode = [0, 0.1, 0, 0]
    elif 500000000 > contract_sp >= 100000000:
        explode = [0, 0, 0.1, 0]
    elif 100000000 > contract_sp:
        explode = [0, 0, 0, 0.1]

    fig = (2.8, 3.3)
    plt.figure(figsize=fig)
    ratio = [over_1G, over_500M, over_100M, under_100M]
    plt.pie(ratio, autopct='%.1f%%', colors=['red', '#00C0AA', '#D1D2D4', '#4C4C4E'], counterclock=False,
            explode=explode)

    image_stream_types = io.BytesIO()
    plt.savefig(image_stream_types, bbox_inches='tight', transparent=True)
    left = Inches(5.8)
    top = Inches(4.9)
    width = Inches(2)
    height = Inches(1.9)
    pic = shapes.add_picture(image_stream_types, left, top, width, height)
    plt.close()


    #####동종업계 IP주소 자원 이용현황#####

    fig = (4, 4)
    plt.figure(figsize=fig)

    ip_list = [1000, 500, 256, 128, 64, 32, 16, 8, 4]
    colors = ['#e5e5e5', '#d4d4d4', '#c7c7c7', '#b6b6b6', '#979797', '#757575', '#595959', '#4c4c4c']
    for index, value in enumerate(ip_list):
        if total_ip >= value:
            print('total_ip: {}'.format(total_ip))
            print('colors: {}\n'.format(colors))
            colors.insert(index, '#00C0AA')
            break

    tick_list = ['1000', '500', '256', '128', '64', '32', '16', '8', '4']
    value = [ip_1000, ip_500, ip_256, ip_128, ip_64, ip_32,
             ip_16, ip_8, ip_4]
    plt.bar(tick_list, value, color=colors)

    image_stream_types = io.BytesIO()
    plt.savefig(image_stream_types, bbox_inches='tight', transparent=True)
    left = Inches(8.1)
    top = Inches(5)
    width = Inches(2.3)
    height = Inches(2.1)
    pic = shapes.add_picture(image_stream_types, left, top, width, height)
    plt.close()

    # img_path = './images/img.png'
    #
    # left = Inches(5.55)
    # top = Inches(1.5)
    # pic = slide.shapes.add_picture(img_path, left, top)

    ############추가 페이지##########

    if int(max(traffic_df['traffic'])) >= int(offer_speed*0.7):
        slide_layout = prs.slide_layouts[4]
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes

    file_name = data1_1 + '_' + data1_3 + '_' + str(datetimes) + '.pptx'
    url = './reports/' + file_name
    prs.save(url)

    return file_name
