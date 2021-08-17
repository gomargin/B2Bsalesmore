from pptx import Presentation
from pptx.util import Inches
from dateutil.relativedelta import relativedelta
import matplotlib
matplotlib.use('Agg')
import matplotlib.font_manager as fm
import sys
import os
import io
import re
import pandas as pd
import numpy as np
import datetime
import pymysql

from zabbix import Zabbix
from model import Scaler, Train
from model import reframe_df

from dateutil import rrule
from matplotlib import pyplot as plt

# tensorflow WARNING 제거
# tf.compat.v1.logging.set_verbosity(tf.compat.v1.logging.ERROR)


class Reporting(Zabbix):
    def __init__(self, row):
        super().__init__()
        """
        Args:
            db_info (type: list)
            - [user, passwd, host, db, port, charset]
            row (type: list)
            - [idx, datetime, line_no, date_start, date_end]
        """
        self.idx = row[0]
        self.datetime = row[2]
        self.line_no = row[3]
        self.date_start = row[4]
        self.date_end = row[5]

    def request(self):
        """웹에서 호출 시 report를 생성하는 함수

        report table의 state column확인:
         - 0 : 리포트 생성 시작
         - 1 : 리포트 생성 중
         - 2 : 리포트 생성 완료
         - 9 : 에러 발생
        """
        self.update(1, self.idx)

        try:
            url = self.generate()
            if url[:5] != 'error':  # error 발생하지 않을 경우
                self.update(2, url, self.idx)
                print('report 생성 완료: {}\n'.format(url))
            else:  # error 발생할 경우
                self.update(9, url, self.idx)
                print('report 생성 중 error 발생', url + '\n', sep='\n')
            self.test_conn.close()

        except:  # 식별되지 않은 예외의 error 발생할 경우
            self.update(9, '식별되지 않은 error', self.idx)
            self.test_conn.close()
            print('report 생성 중 식별되지 않은 error 발생\n')

    def generate(self):
        print('회선번호: {} report 생성 시작...\n'.format(self.line_no))
        line_num = self.line_no.replace('-', '')

        # neoss table 읽기
        sql = 'SELECT service_id, host_name, interface, ethernet_ip FROM neoss WHERE leased_line_num="{}"'.format(
            line_num)
        neoss_df = self.read(sql)
        if not neoss_df.empty:
            neoss_df.columns = ['service_id', 'host_name', 'interface', 'ethernet_ip']
            print('----- neoss_df -----')
            print(neoss_df, '\n')
        else:
            error = 'error: More than one column of neoss_df is empty'
            return error

        # neoss_df 에서 ip 데이터 뽑기. ppt #3  'ip 할당 현황' 데이터
        ethernet_ip = neoss_df['ethernet_ip'].values

        ip_count = [0 for _ in range(7)]  # ip count 하는 list [256, 128, 64, 32, 16, 8, 4] 순서
        total_ip = 0
        for ip in ethernet_ip:
            ip_end = int(ip[-2:])  # ip 마지막 두 자리 정수로 변경
            ip_count[ip_end - 24] += 1  # ip 마지막 두 자리 정수 변환. 24 -> 256, 25 -> 128, 24 -> 64 ...
            total_ip += pow(2, 32 - ip_end)

        ethernet_ip_24 = ip_count[0]
        ethernet_ip_25 = ip_count[1]
        ethernet_ip_26 = ip_count[2]
        ethernet_ip_27 = ip_count[3]
        ethernet_ip_28 = ip_count[4]
        ethernet_ip_29 = ip_count[5]
        ethernet_ip_30 = ip_count[6]

        ethernet_ip_list = ['-' for i in range(8)]
        for i in range(len(ethernet_ip)):
            if i < 8:
                start_ip = int(ethernet_ip[i].split('.')[-1].split('/')[0])
                ip_num = pow(2, 32 - int(ethernet_ip[i][-2:]))
                ethernet_ip_list[i] = ethernet_ip[i][:-3] + '~' + str(start_ip + ip_num - 1)
            else:
                break

        # neoss_df 에서 service_id 값 뽑기
        service_id = neoss_df['service_id'][0]

        # bidw table 에서 데이터 가져오기
        sql = 'SELECT industry_2, customer_name, contract_speed, service_name, contract_end FROM bidw WHERE ' \
                      'service_id="{}"'.format(service_id)
        bidw_df = self.read(sql)
        if not bidw_df.empty:
            bidw_df.columns = ['industry_2', 'customer_name', 'contract_speed', 'service_name', 'contract_end']
            print('----- bidw_df -----')
            print(bidw_df, '\n')
        else:
            error = 'error: More than one column of bidw_df is empty'
            return error

        # bidw table 에서 industry 기준 contract_speed
        industry = bidw_df['industry_2'][0]
        sql = 'SELECT industry_2, contract_speed FROM bidw WHERE industry_2="{}"'.format(industry)
        industry_speed_df = self.read(sql)
        if not industry_speed_df.empty:
            industry_speed_df.columns = ['industry_2', 'contract_speed']
            print('----- industry_speed_df -----')
            print(industry_speed_df.head(), '\n')
        else:
            error = 'error: More than one column of industry_speed_df is empty'
            return error

        # contract_speed 값 변환 (M -> pow(10,6), G -> pow(10,9))
        speed_data = []
        for i in industry_speed_df['contract_speed']:
            if i[-1] == 'M':
                speed_data.append(int(i[:-1]) * pow(10, 6))  # 수정
            elif i[-1] == 'G':
                speed_data.append(int(i[:-1]) * pow(10, 9))  # 수정

        # speed 개수 count
        over_1G = 0
        over_500M = 0
        over_100M = 0
        under_100M = 0
        for i in speed_data:  # 수정
            if i >= pow(10, 9):
                over_1G += 1
            elif pow(10, 9) > i >= 5 * pow(10, 8):
                over_500M += 1
            elif 5 * pow(10, 8) > i >= pow(10, 8):
                over_100M += 1
            else:
                under_100M += 1
        print('----- ip 개수 -----')
        print("1G 이상: {} 개".format(over_1G))
        print("500M 이상: {} 개".format(over_500M))
        print("100M 이상: {} 개".format(over_100M))
        print("100M 미만: {} 개".format(under_100M), '\n')

        # bidw table 과 neoss table 에서 service_id 기준 industry, customer_name, ethernet_ip 합치기. 동종업계 ip 비교 위함
        sql = 'SELECT industry_2, customer_name, service_id FROM bidw WHERE industry_2="{}"'.format(industry)
        merge_df_1 = self.read(sql)
        if not merge_df_1.empty:
            merge_df_1.columns = ['industry_2', 'customer_name', 'service_id']
            print('----- merge_df_1 -----')
            print(merge_df_1.head(), '\n')
        else:
            error = 'error: More than one column of merge_df_1 is empty'
            return error

        sql = 'SELECT service_id, ethernet_ip FROM neoss'
        merge_df_2 = self.read(sql)
        if not merge_df_2.empty:
            merge_df_2.columns = ['service_id', 'ethernet_ip']
            print('----- merge_df_2 -----')
            print(merge_df_2.head(), '\n')
        else:
            error = 'error: More than one column of merge_df_2 is empty'
            return error

        merge_df = pd.merge(merge_df_1, merge_df_2, on="service_id")
        print('----- merge_df -----')
        print(merge_df.head(), '\n')

        # 동종업계 고객 별 ip value 설정
        customer_name_dict = {}
        for customer in merge_df['customer_name']:
            customer_name_dict[customer] = 0

        for row in merge_df.values:  # col1: industry, col2: customer_name, col3: service_id, col4: ethernet_ip
            if row[3] is not None:
                customer_name_dict[row[1]] += pow(2, 32 - int(row[3][-2:]))
            else:
                pass

        # 동종업계 ip 별 count
        ip_list = [1000, 500, 256, 128, 64, 32, 16, 8, 4]
        customer_ip_count = [0 for _ in range(9)]

        for ip in customer_name_dict.values():
            for index, value in enumerate(ip_list):
                if ip >= ip_list[index]:
                    customer_ip_count[index] += 1
                    break
        ip_1000 = customer_ip_count[0]
        ip_500 = customer_ip_count[1]
        ip_256 = customer_ip_count[2]
        ip_128 = customer_ip_count[3]
        ip_64 = customer_ip_count[4]
        ip_32 = customer_ip_count[5]
        ip_16 = customer_ip_count[6]
        ip_8 = customer_ip_count[7]
        ip_4 = customer_ip_count[8]

        print('----- customer_ip_count -----')
        print(customer_ip_count, '\n')

        # host_name 과 hostid 추출
        print('----- host info -----')
        host_name = neoss_df['host_name'][0]
        print('host_name: {}'.format(host_name))

        sql = 'SELECT hostid FROM hosts WHERE host="{}"'.format(host_name)
        hostid = self.read_one(sql)
        hostid = hostid[0]
        print('hostid: {}'.format(hostid), '\n')

        # itemid 추출
        print('----- interface & itemid -----')
        interface = neoss_df['interface'][0]
        print('interface: {}'.format(interface))
        # key_값 검색을 위한 정규식 표현 생성
        regexp_list= []
        # 1. 대다수의 일반적인 interface
        regexp_list.append('.+(' + interface + ')[^0-9]')
        # 2. interface 에 ifHC~ 포함된 경우
        if ('ifHCInOctets' in interface) or ('ifHCOutOctets' in interface):
            regexp_list.append('.+(' + interface[interface.find('[')+1:-1] + ')[^0-9]')
        # 3. 'ge3/1/19, IPinterface' 에서 숫자를 추출해야 하는 경우
        try:
            p = re.compile('[0-9]+(/)[0-9]+(/)[0-9]+')
            interface_num = p.search(interface).group()
            regexp_list.append('.+(' + interface_num + ')[^0-9]')
        except AttributeError:
            pass

        print('regexp_list: {}'.format(regexp_list))
        for index, regexp in enumerate(regexp_list):
            sql = 'SELECT itemid, key_ FROM items WHERE hostid="{}" and key_ REGEXP "{}"'.format(str(hostid), regexp)
            itemid_df = self.read(sql)
            if not itemid_df.empty:
                itemid_df.columns = ['itemid', 'key_']
                if 'In' in itemid_df['key_'][0]:
                    in_itemid, out_itemid = itemid_df['itemid']
                else:
                    out_itemid, in_itemid = itemid_df['itemid']
                itemid_list = [in_itemid, out_itemid]
                print('regexp: {}'.format(regexp))
                print('in_itemid: {}, out_itemid: {}'.format(in_itemid, out_itemid), '\n')
                break
            elif index == len(regexp_list)-1:
                print('interface: {}'.format(interface))
                print('regexp_list: {}'.format(regexp_list))
                error = 'error: interface 값을 통한 itemid 찾기 실패'
                return error

        # 제공속도와 청약속도 추출
        sql = 'SELECT engre, gbic FROM realspeed WHERE leased_line_num={}'.format(line_num)
        system_df = self.read(sql)
        if not system_df.empty:
            system_df.columns = ['engre', 'gbic']
            print('----- system_df -----')
            print(system_df, '\n')
        else:
            error = 'error: More than one column of system_df is empty'
            return error

        # offer_speed: 제공속도, contract_speed: 청약속도, _M: 각 속도 메가 단위
        offer_speed = int(system_df['engre'][0])
        offer_speed_M = int(offer_speed/pow(10, 6))
        contract_speed = int(system_df['gbic'][0])
        contract_speed_M = int(contract_speed/pow(10, 6))

        if int(offer_speed) >= contract_speed:
            if offer_speed >= 1000*pow(10, 6):
                script2_29 = '시설 변경 필요'
            else:
                script2_29 = '증속 불가'
            script2_30 = '중단'
        else:
            script2_29 = '시설변경 없이\n' \
                         '추가 ' + str(contract_speed_M - offer_speed_M) + 'M 증속 가능'
            script2_30 = '무중단 작업 가능'


        # AI 학습 및 예측
        for itemid in itemid_list:
            sql = "SELECT itemid, from_unixtime(clock) as datetime, value as traffic " \
                          "FROM history_uint WHERE itemid IN ({})".format(itemid)
            df = self.read(sql)
            if not df.empty:
                df.columns = ['itemid', 'datetime', 'traffic']
                print('----- itemid df -----')
                print(df.head(), '\n')
            else:
                error = 'error: More than one column of df is empty'
                return error

            dp_df = reframe_df(df)
            print('----- dp_df -----')
            print(dp_df, '\n')
            traffic_scaler = Scaler(dp_df, 'traffic')
            sc_df = traffic_scaler.normalization()
            df_len = len(sc_df)
            print('df_len: {}'.format(df_len))
            if df_len >= 42:
                pred_len = df_len // 6
            else:
                pred_len = 0
            print('pred_len: {}'.format(pred_len))

            in_steps = 3
            out_steps = 1

            # 해당 itemid에 해당하는 df 생성

            if pred_len != 0:
                upload_df = pd.DataFrame(index=range(0, pred_len), columns=['itemid', 'dates', 'traffic'])
                upload_df['itemid'] = itemid
                # 예측하고자 하는 첫날부터 out_steps 만큼의 list 생성
                upload_df['dates'] = [(sc_df.index[-1] + datetime.timedelta(days=i)).strftime('%y%m%d') for i in range(pred_len)]
                trainer = Train(itemid, pred_len, in_steps, out_steps, valid_per=0, epochs=100, batch_size=128, unit=128, drop_per=0.1)
                trainer.train_model(sc_df)
                # 음의 예측값은 0으로 변환
                real_prediction = traffic_scaler.rev_normalization(trainer.predict_model(sc_df))
                upload_df['traffic'] = [max(num, 0) for num in real_prediction]
                print('예측 완료!', '\n')
            else:
                upload_df = pd.DataFrame(index=range(7), columns=['itemid', 'dates', 'traffic'])
                upload_df['itemid'] = itemid
                upload_df['dates'] = [(sc_df.index[-1] + datetime.timedelta(days=i)).strftime('%y%m%d') for i in range(7)]
                upload_df['traffic'] = 0
                print('데이터 부족으로 예측 불가', '\n')
        self.delete(itemid)
        self.insert(upload_df)

        # 실제 traffic 값 추출
        real_traffic_list = []
        for itemid in itemid_list:
            sql = 'SELECT itemid, from_unixtime(clock), value  FROM history_uint WHERE itemid="{}" ' \
                             'and "{}" <= from_unixtime(clock) AND from_unixtime(clock) <= "{}"'\
                            .format(itemid, self.date_start, self.date_end)
            traffic_df = self.read(sql)
            if not traffic_df.empty:
                traffic_df.columns = ['itemid', 'clock', 'traffic']
                real_traffic_list.append(traffic_df)
            else:
                error = 'error: More than one column of traffic_df is empty'
                return error

        in_real_traffic_df = real_traffic_list[0]
        out_real_traffic_df = real_traffic_list[1]
        print('-----in_real_traffic_df-----')
        print(in_real_traffic_df, '\n')
        print('-----out_real_traffic_df-----')
        print(out_real_traffic_df, '\n')


        # 트래픽 임계치 count 행렬
        # 1행: in, 2행: out
        # 0열~6열: 제공속도 100% 초과/ 90~100% / 80~90% / 70~80% / 50~70% / 49% 이하
        traffic_count = np.zeros((2, 6), dtype='int')

        for index, df in enumerate(real_traffic_list):
            for value in df['traffic']:
                try:
                    if int(value) > int(offer_speed):
                        traffic_count[index, 0] += 1
                    elif int(offer_speed) > int(value) >= int(offer_speed)*0.9:
                        traffic_count[index, 1] += 1
                    elif int(offer_speed)*0.9 > int(value) >= int(offer_speed)*0.7:
                        traffic_count[index, 2] += 1
                    elif int(offer_speed)*0.7 > int(value) >= int(offer_speed)*0.01:
                        traffic_count[index, 3] += 1
                    else:
                        traffic_count[index, 4] += 1
                except:
                    print('try 에러 발생!')
                    pass

        in_over_100 = traffic_count[0, 0]
        in_over_90 = traffic_count[0, 1]
        in_over_70 = traffic_count[0, 2]
        in_under_69 = traffic_count[0, 3]
        in_under_1 = traffic_count[0, 4]
        out_over_100 = traffic_count[1, 0]
        out_over_90 = traffic_count[1, 1]
        out_over_70 = traffic_count[1, 2]
        out_under_69 = traffic_count[1, 3]
        out_under_1 = traffic_count[1, 4]

        arr_list = []
        for row in range(traffic_count.shape[0]):
            for col in range(traffic_count.shape[1]):
                if traffic_count[row][col] != 0:
                    arr_list.append((row, col))
                    break
        if (arr_list[0][1] <= 3) and (arr_list[1][1] <= 3):   # 업/다운 둘 중 하나라도 70% 넘을 경우
            if arr_list[0][1] == arr_list[1][1]:   # % 동일
                updown = '업/다운로드'
                max_col = arr_list[0][1]
                over_count = in_over_100 + out_over_100
        else:
            max_col = max(arr_list[0][1], arr_list[1][1])
            if arr_list[0][1] > arr_list[1][1]:   # 업로드가 더 큰 경우
                updown = '업로드'
                over_count = traffic_count[0, max_col]
            else:   # 다운로드가 더 큰 경우
                updown = '다운로드'
                over_count = traffic_count[1, max_col]

        if max_col == 0:
            script1_12_0 = '{} 트래픽이 제공속도인 {}Mbps를 초과하는 트래픽 폭주가 {}회 발생하였습니다.'.format(updown, offer_speed_M, over_count)
            script1_12_1 = '트래픽 폭주가 지속적으로 발생하는 경우, 접속불량, 속도지연 등 서비스 이용에 불편을 초래할 수 있습니다.'
            script1_12 = [script1_12_0, script1_12_1]

            script1_14_0 = '정상적인 업무 트래픽에 의한 폭주일 경우, GiGA Office 도입 또는 코넷 증속을 권고합니다.'
            script1_14_1 = '보다 안정적인 서비스 이용을 위해 KT만의 차별화된 맞춤형 보안서비스를 이용할 수 있습니다.'
            script1_14_2 = '보안진단 컨설팅(E-브로셔 참조) 이용 및 보안서비스 가입은 영업 담당자에게 문의하여 주시기 바랍니다.'
            script1_14 = [script1_14_0, script1_14_1, script1_14_2]
        elif max_col == 1:
            script1_12_0 = '{} 트래픽이 제공속도인 {}Mbps를 초과하는 트래픽 폭주가 {}회 발생하였습니다.'.format(updown, offer_speed_M, over_count)
            script1_12_1 = '단, 최고치 트래픽이 5분 평균치로 산출되므로, 실제 체감속도는 더욱 낮을 것으로 판단됩니다.'
            script1_12 = [script1_12_0, script1_12_1]

            script1_14_0 = '정상적인 업무 트래픽에 의한 폭주일 경우, GiGA Office 도입 또는 코넷 증속을 권고합니다.'
            script1_14_1 = '보다 안정적인 서비스 이용을 위해 KT만의 차별화된 맞춤형 보안서비스를 이용할 수 있습니다.'
            script1_14_2 = '보안진단 컨설팅(E-브로셔 참조) 이용 및 보안서비스 가입은 영업 담당자에게 문의하여 주시기 바랍니다.'
            script1_14 = [script1_14_0, script1_14_1, script1_14_2]
        elif max_col == 2:
            script1_12_0 = '{} 트래픽의 안정적 이용 권고치 70%를 초과한 건수가 {}회 발생하였습니다.'.format(updown, over_count)
            script1_12_1 = '서비스 이용 상 문제는 발생하지 않을 것으로 판단되지만, 지속적인 관리가 필요합니다.'
            script1_12 = [script1_12_0, script1_12_1]

            script1_14_0 = '지속적인 트래픽 증가 트렌드를 감안하여 GiGA Office 도입 또는 코넷 증속을 권고합니다.'
            script1_14_1 = '보다 안정적인 서비스 이용을 위해 KT만의 차별화된 맞춤형 보안서비스를 이용할 수 있습니다.'
            script1_14_2 = '보안진단 컨설팅(E-브로셔 참조) 이용 및 보안서비스 가입은 영업 담당자에게 문의하여 주시기 바랍니다.'
            script1_14 = [script1_14_0, script1_14_1, script1_14_2]
        elif max_col == 3:
            script1_12_0 = '업/다운로드 양방향 모두, 제공속도의 70% 권고치 이내에서 안정적으로 이용하고 있습니다.'
            script1_12_1 = '서비스 이용 상 문제는 발생하지 않을 것으로 판단되지만, 부하증가 요인을 감안하여 지속적인 관리가 필요합니다.'
            script1_12 = [script1_12_0, script1_12_1]

            script1_14_0 = '지속적인 트래픽 증가 트렌드를 감안하여 주기적인 트래픽 모니터링 관리를 권고합니다.'
            script1_14_1 = '보다 안정적인 서비스 이용을 원하시면, KT만의 차별화된 맞춤형 보안서비스를 이용할 수 있습니다.'
            script1_14_2 = '자세한 사항은 영업 담당자에게 문의하여 주시기 바랍니다.'
            script1_14 = [script1_14_0, script1_14_1, script1_14_2]
        elif max_col == 4:
            script1_12_0 = '업/다운로드 양방향 모두, 트래픽이 측정되지 않았습니다.'
            script1_12_1 = '실제 사용 중인 경우, 트래픽이 수집되지 않았을 수 있으니, KT 담당자에게 문의 바랍니다.'
            script1_12 = [script1_12_0, script1_12_1]

            script1_14_0 = '트래픽이 수집되지 않고 있으며, 세부 원인은 추가 확인이 필요합니다.'
            script1_14_1 = '정상적으로 서비스를 이용하고 있는 경우, 영업 담당자에게 문의하여 주시기 바랍니다.'
            script1_14_2 = '평소 이용하지 않는 무휴/예비 회선인 경우, 비상시를 대비한 주기적인 관리가 필요합니다.'
            script1_14 = [script1_14_0, script1_14_1, script1_14_2]



        # AI 예측 데이터 추출
        pred_traffic_list = []
        for itemid in itemid_list:
            sql = 'SELECT itemid, dates, traffic FROM test_ys WHERE itemid="{}"'.format(itemid)
            pred_traffic_df = self.read(sql)
            if not pred_traffic_df.empty:
                pred_traffic_df.columns = ['itemid', 'dates', 'traffic']
                pred_traffic_list.append(pred_traffic_df)
            else:
                error = 'error: More than one column of pred_traffic_df is empty'
                return error

        in_pred_traffic_df = pred_traffic_list[0]
        out_pred_traffic_df = pred_traffic_list[1]
        print('----- in_pred_traffic_df -----')
        print(in_pred_traffic_df, '\n')
        print('----- out_pred_traffic_df -----')
        print(out_pred_traffic_df, '\n')

        dates = ['20' + i[:2] + '-' + i[2:4] + '-' + i[4:] for i in pred_traffic_df['dates']]

        # 예측치와 제공속도 비교
        pred_over_count = np.zeros((2, 3), dtype='int')

        for index, df in enumerate(pred_traffic_list):
            for traffic in df['traffic']:
                try:
                    if int(traffic) > int(offer_speed)*0.9:
                        pred_over_count[index, 0] += 1
                    elif int(offer_speed)*0.9 > int(traffic) >= int(offer_speed)*0.7:
                        pred_over_count[index, 1] += 1
                    else:
                        pred_over_count[index, 2] += 1
                except:
                    print('try 에러 발생!')
                    pass

        pred_in_over_90 = pred_over_count[0, 0]
        pred_in_over_70 = pred_over_count[0, 1]
        pred_in_under_70 = pred_over_count[0, 2]
        pred_out_over_90 = pred_over_count[1, 0]
        pred_out_over_70 = pred_over_count[1, 1]
        pred_out_under_70 = pred_over_count[1, 2]

        if (pred_in_over_90 + pred_out_over_90) > 0:
            script1_13_0 = 'AI 예측 알고리즘을 통한 분석결과 지속적인 폭주가 예상 됨'
            script1_13_1 = '향후 권고치' + str(int(offer_speed_M*0.7)) + 'M를 초과 할 것으로 예측 됨'
            script1_13 = [script1_13_0, script1_13_1]
        elif (pred_in_over_70 + pred_out_over_70) > 0:
            script1_13_0 = 'AI 예측 알고리즘을 통한분석결과 주기적 폭주가 예상 됨'
            script1_13 = [script1_13_0]
        else:
            script1_13_0 = 'AI 예측 알고리즘을 통한 분석결과 완만한 증가가 예상 됨'
            script1_13 = [script1_13_0]


        #########total script###########
        print('report 생성을 위한 데이터 정리 시작...')

        # -------------------한글 폰트 지원-------------------
        path_gothic = 'C:/Windows/Fonts/malgunbd.ttf'
        fontprop1 = fm.FontProperties(fname=path_gothic, size=10)

        # -----------------------페이지별 항목----------------------------
        prs_test = Presentation('Sales_More_ppt2.2.pptx')
        slide0 = prs_test.slides.add_slide(prs_test.slide_layouts[0])
        # for shape in slide0.placeholders:
        #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
        # slide1 = prs_test.slides.add_slide(prs_test.slide_layouts[1])
        # for shape in slide1.placeholders:
        #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
        # slide2 = prs_test.slides.add_slide(prs_test.slide_layouts[2])
        # for shape in slide2.placeholders:
        #     print('%d %s' % (shape.placeholder_format.idx, shape.name))


        data0_1 = bidw_df['customer_name'][0]
        data0_2 = str(datetime.datetime.now().date())
        print('첫 번째 슬라이드 데이터 정의 완료')
        
        data1_1 = bidw_df['customer_name'][0]
        data1_2 = neoss_df['service_id'][0]   # 수정: df_neoss -> neoss_df
        data1_3 = line_num[:2] + '-' + line_num[2:8] + '-' + line_num[-4:]
        data1_4 = str(offer_speed_M) + 'M'
        data1_5 = self.date_start
        data1_6 = self.date_end
        data1_7 = dates[0]
        data1_8 = dates[-1]
        data1_9 = str(int(in_real_traffic_df['traffic'].max()/pow(10, 6))) + 'M / ' + str(int(out_real_traffic_df['traffic'].max()/pow(10, 6))) + 'M'
        data1_10 = format((in_real_traffic_df['traffic'].max()/offer_speed)*100, '.1f') + '% / ' + format((out_real_traffic_df['traffic'].max()/offer_speed)*100, '.1f') + '%'
        data1_11 = str(int(in_pred_traffic_df['traffic'].max()/pow(10, 6))) + 'M / ' + str(int(out_pred_traffic_df['traffic'].max()/pow(10, 6))) + 'M'
        data1_12 = format((in_pred_traffic_df['traffic'].max()/offer_speed)*100, '.1f') + '% / ' + format((out_pred_traffic_df['traffic'].max()/offer_speed)*100, '.1f') + '%'
        print('두 번째 슬라이드 데이터 정의 완료')
        
        # 약정만료 기한 표시
        contract_end = bidw_df['contract_end'][0]
        if contract_end == '999912':
            contract_end = '999911'
        contract_ends = datetime.datetime.strptime(contract_end, '%Y%m') + relativedelta(months=1) - datetime.timedelta(days=1)
        now = datetime.datetime.strptime(self.datetime, '%Y-%m-%d-%H%M%S')

        data2_1 = bidw_df['service_name'][0]
        # data2_2 = bidw_df['contract_speed'][0]   청약속도
        data2_3 = str(offer_speed_M) + 'M'
        if contract_ends > now:
            data2_4 = contract_end[:4] + '년 ' +  contract_end[4:] + '월'
        else:
            data2_4 = '-'
        data2_5 = str(ethernet_ip_24)
        data2_6 = str(ethernet_ip_25)
        data2_7 = str(ethernet_ip_26)
        data2_8 = str(ethernet_ip_27)
        data2_9 = str(ethernet_ip_28)
        data2_10 = str(ethernet_ip_29)
        data2_11 = ethernet_ip_list[0]
        data2_12 = ethernet_ip_list[1]
        data2_13 = ethernet_ip_list[2]
        data2_14 = ethernet_ip_list[3]
        data2_15 = ethernet_ip_list[4]
        data2_16 = ethernet_ip_list[5]
        data2_17 = ethernet_ip_list[6]
        data2_18 = ethernet_ip_list[7]
        data2_19 = str(in_over_70 + out_over_70)
        data2_20 = str(in_over_90 + out_over_90)
        data2_21 = str(in_over_100 + out_over_100)
        data2_22 = str(int(data2_19) + int(data2_20) + int(data2_21))
        data2_23 = '2_23'
        data2_24 = '2_24'
        data2_25 = '2_25'
        data2_26 = '2_26'
        data2_27 = str(offer_speed_M) + 'M'
        data2_28 = str(contract_speed_M) + 'M'
        data2_29 = script2_29
        data2_30 = script2_30
        data2_31 = industry
        data2_32 = str(ethernet_ip_30)
        data2_33 = '2_33'
        data2_34 = '2_34'
        data2_35 = '2_35'
        data2_36 = '2_36'
        data2_37 = '2_37'
        data2_38 = '2_38'
        data2_39 = '2_39'
        data2_40 = '2_40'
        print('세 번째 슬라이드 데이터 정의 완료')

        # ppt 파일 불러오기
        prs = Presentation('Sales_More_ppt2.2.pptx')
        # 슬라이드 마스터 첫번째 레이아웃 불러오기 (prs.slide_layouts)
        slide_layout = prs.slide_layouts[0]
        # 슬라이드 생성하기
        slide = prs.slides.add_slide(slide_layout)
        # 슬라이드 마스터의 개체 틀
        shapes = slide.shapes

        #####page0####
        page0_1_frame = shapes.placeholders[20].text_frame
        page0_1_frame.text = data0_1
        page0_2_frame = shapes.placeholders[19].text_frame
        page0_2_frame.text = data0_2
        print('첫 번째 슬라이드 생성 완료')

        #####page1####
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes

        page1_1_frame = shapes.placeholders[14].text_frame
        page1_1_frame.text = data1_1
        page1_2_frame = shapes.placeholders[15].text_frame
        page1_2_frame.text = data1_2
        page1_3_frame = shapes.placeholders[16].text_frame
        page1_3_frame.text = data1_3
        page1_4_frame = shapes.placeholders[17].text_frame
        page1_4_frame.text = data1_4
        page1_5_frame = shapes.placeholders[29].text_frame
        page1_5_frame.text = data1_5
        page1_6_frame = shapes.placeholders[30].text_frame
        page1_6_frame.text = data1_6
        page1_7_frame = shapes.placeholders[31].text_frame
        page1_7_frame.text = data1_7
        page1_8_frame = shapes.placeholders[32].text_frame
        page1_8_frame.text = data1_8
        page1_9_frame = shapes.placeholders[18].text_frame
        page1_9_frame.text = data1_9
        page1_10_frame = shapes.placeholders[27].text_frame
        page1_10_frame.text = data1_10
        page1_11_frame = shapes.placeholders[33].text_frame
        page1_11_frame.text = data1_11
        page1_12_frame = shapes.placeholders[34].text_frame
        page1_12_frame.text = data1_12

        if max_col == 0:
            page1_12_frame = shapes.placeholders[22].text_frame.paragraphs[0]
            page1_12_frame.text = script1_12[0]
            page1_12_frame_1 = shapes.placeholders[22].text_frame.add_paragraph()
            page1_12_frame_1.text = script1_12[1]
        else:
            page1_12_frame = shapes.placeholders[22].text_frame.paragraphs[0]
            page1_12_frame.text = script1_12[0]
            page1_12_frame_1 = shapes.placeholders[22].text_frame.add_paragraph()
            page1_12_frame_1.text = script1_12[1]
            page1_12_frame_1.level = 1

        if len(script1_13) == 2:
            page1_13_frame = shapes.placeholders[25].text_frame.paragraphs[0]
            page1_13_frame.text = script1_13[0]
            page1_13_frame_1 = shapes.placeholders[25].text_frame.add_paragraph()
            page1_13_frame_1.text = script1_13[1]
        else:
            page1_13_frame = shapes.placeholders[25].text_frame.paragraphs[0]
            page1_13_frame.text = script1_13[0]

        if max_col == 4:
            page1_14_frame = shapes.placeholders[26].text_frame.paragraphs[0]
            page1_14_frame.text = script1_14[0]
            page1_14_frame_1 = shapes.placeholders[26].text_frame.add_paragraph()
            page1_14_frame_1.text = script1_14[1]
            page1_14_frame_1.level = 1
            page1_14_frame_2 = shapes.placeholders[26].text_frame.add_paragraph()
            page1_14_frame_2.text = script1_14[2]
            page1_14_frame_2.level = 1
        else:
            page1_14_frame = shapes.placeholders[26].text_frame.paragraphs[0]
            page1_14_frame.text = script1_14[0]
            page1_14_frame_1 = shapes.placeholders[26].text_frame.add_paragraph()
            page1_14_frame_1.text = script1_14[1]
            page1_14_frame_2 = shapes.placeholders[26].text_frame.add_paragraph()
            page1_14_frame_2.text = script1_14[2]


        ##### Traffic plot #####

        fig = (14, 2)
        plt.figure(figsize=fig)

        plt.title('M', loc="left", fontsize="10")
        plt.plot(in_real_traffic_df['traffic']/pow(10, 6), linestyle='-.', linewidth=1, color='blue', label='Up')
        plt.plot(out_real_traffic_df['traffic']/pow(10, 6), linestyle=':', linewidth=1, color='red', label='Down')
        plt.legend(loc='upper right')

        xticks_value = [i for i in range(0, len(in_real_traffic_df), (len(in_real_traffic_df) // 10))]
        xticks = [in_real_traffic_df['clock'][i].strftime('%Y%m%d%H%M')[2:8] for i in xticks_value]
        plt.xticks(xticks_value, xticks, rotation=45, fontsize=12)

        y_lim = max(max(in_real_traffic_df['traffic']), max(out_real_traffic_df['traffic'])) * 1.2
        plt.ylim(0, y_lim / pow(10, 6))

        if int(offer_speed) * 0.7 <= max(out_real_traffic_df['traffic']) * 1.2 < int(offer_speed):
            plt.axhline(y=int(offer_speed_M * 0.7), color='orange', linewidth=1, linestyle='--')
            plt.text(0, int(offer_speed_M * 0.7), "권고속도 : " + str(int(offer_speed_M * 0.7)) + 'M', fontsize=10, fontproperties=fontprop1)
            # plt.text(0, int(offer_speed_M * 0.7), "권고속도 : " + str(int(offer_speed_M * 0.7)) + 'M', fontsize=10)
        elif int(offer_speed) <= max(out_real_traffic_df['traffic']) * 1.2:
            plt.axhline(y=offer_speed_M, color='red', linewidth=1, linestyle='--')
            plt.text(0, offer_speed_M, "제공속도 : " + str(offer_speed_M) + 'M', fontsize=10, fontproperties=fontprop1)
            # plt.text(0, offer_speed_M, "제공속도 : " + str(offer_speed_M) + 'M', fontsize=10)
            plt.axhline(y=int(offer_speed_M * 0.7), color='orange', linewidth=1, linestyle='--')
            plt.text(0, int(offer_speed_M * 0.7), "권고속도 : " + str(int(offer_speed_M * 0.7)) + 'M', fontsize=10, fontproperties=fontprop1)
            # plt.text(0, int(offer_speed_M * 0.7), "권고속도 : " + str(int(offer_speed_M * 0.7)) + 'M', fontsize=10)

        image_stream_types = io.BytesIO()
        plt.savefig(image_stream_types, bbox_inches='tight', transparent=True)
        left = Inches(0.35)
        top = Inches(2.3)
        width = Inches(6)
        height = Inches(1.9)
        pic = shapes.add_picture(image_stream_types, left, top, width, height)
        plt.close()


        ##### AI plot #####

        if pred_len == 0:   # 데이터 부족 시 그림 삽입
            data_img_path = './images/data_lack.jpg'
            left_cm = 17.93
            left = Inches(left_cm / 2.54)  # inch = cm/2.54
            top_cm = 6.13
            top = Inches(top_cm / 2.54)
            pic = slide.shapes.add_picture(data_img_path, left, top)
        else:
            fig = (8, 2)
            plt.figure(figsize=fig)

            plt.title('M', loc="left", fontsize="10")
            plt.plot(in_pred_traffic_df['traffic']/pow(10, 6), linestyle='-.', linewidth=1, color='blue', label='Up')
            plt.plot(out_pred_traffic_df['traffic']/pow(10, 6), linestyle=':', linewidth=1, color='red', label='Down')
            plt.legend(loc='upper right')

            xticks_value = [i for i in range(0, len(in_pred_traffic_df), 1)]
            xticks = [in_pred_traffic_df['dates'][i] for i in xticks_value]
            plt.xticks(xticks_value, xticks, rotation=45, fontsize=12)

            plt.ylim(0, y_lim / pow(10, 6))   # 예측 그래프의 y축 limit은 실제 그래프와 동일

            if int(offer_speed) * 0.7 <= max(out_real_traffic_df['traffic']) * 1.2 < int(offer_speed):
                plt.axhline(y=int(offer_speed_M * 0.7), color='orange', linewidth=1, linestyle='--')
                plt.text(0, int(offer_speed_M * 0.7), "권고속도 : " + str(int(offer_speed_M * 0.7)) + 'M', fontsize=10, fontproperties=fontprop1)
                # plt.text(0, int(offer_speed_M * 0.7), "권고속도 : " + str(int(offer_speed_M * 0.7)) + 'M', fontsize=10)
            elif int(offer_speed) <= max(out_real_traffic_df['traffic']) * 1.2:
                plt.axhline(y=offer_speed_M, color='red', linewidth=1, linestyle='--')
                plt.text(0, offer_speed_M, "제공속도 : " + str(offer_speed_M) + 'M', fontsize=10, fontproperties=fontprop1)
                # plt.text(0, offer_speed_M, "제공속도 : " + str(offer_speed_M) + 'M', fontsize=10)
                plt.axhline(y=int(offer_speed_M * 0.7), color='orange', linewidth=1, linestyle='--')
                plt.text(0, int(offer_speed_M * 0.7), "권고속도 : " + str(int(offer_speed_M * 0.7)) + 'M', fontsize=10, fontproperties=fontprop1)
                # plt.text(0, int(offer_speed_M * 0.7), "권고속도 : " + str(int(offer_speed_M * 0.7)) + 'M', fontsize=10)

            image_stream_types = io.BytesIO()
            plt.savefig(image_stream_types, bbox_inches='tight', transparent=True)
            left = Inches(6.9)
            top = Inches(2.3)
            width = Inches(3.5)
            height = Inches(1.9)
            pic = shapes.add_picture(image_stream_types, left, top, width, height)
            plt.close()
        print('두 번째 슬라이드 생성 완료')

        #####page2#####
        slide_layout = prs.slide_layouts[2]
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes

        page2_1_frame = shapes.placeholders[29].text_frame
        page2_1_frame.text = data2_1
        # page2_2_frame = shapes.placeholders[33].text_frame
        # page2_2_frame.text = data2_2
        page2_3_frame = shapes.placeholders[34].text_frame
        page2_3_frame.text = data2_3
        page2_4_frame = shapes.placeholders[39].text_frame
        page2_4_frame.text = data2_4
        page2_5_frame = shapes.placeholders[40].text_frame
        page2_5_frame.text = data2_5
        page2_6_frame = shapes.placeholders[41].text_frame
        page2_6_frame.text = data2_6
        page2_7_frame = shapes.placeholders[42].text_frame
        page2_7_frame.text = data2_7
        page2_8_frame = shapes.placeholders[43].text_frame
        page2_8_frame.text = data2_8
        page2_9_frame = shapes.placeholders[44].text_frame
        page2_9_frame.text = data2_9
        page2_10_frame = shapes.placeholders[45].text_frame
        page2_10_frame.text = data2_10
        page2_11_frame = shapes.placeholders[73].text_frame
        page2_11_frame.text = data2_11
        page2_12_frame = shapes.placeholders[74].text_frame
        page2_12_frame.text = data2_12
        page2_13_frame = shapes.placeholders[75].text_frame
        page2_13_frame.text = data2_13
        page2_14_frame = shapes.placeholders[76].text_frame
        page2_14_frame.text = data2_14
        page2_15_frame = shapes.placeholders[35].text_frame
        page2_15_frame.text = data2_15
        page2_16_frame = shapes.placeholders[36].text_frame
        page2_16_frame.text = data2_16
        page2_17_frame = shapes.placeholders[37].text_frame
        page2_17_frame.text = data2_17
        page2_18_frame = shapes.placeholders[72].text_frame
        page2_18_frame.text = data2_18
        page2_19_frame = shapes.placeholders[52].text_frame
        page2_19_frame.text = data2_19
        page2_20_frame = shapes.placeholders[53].text_frame
        page2_20_frame.text = data2_20
        page2_21_frame = shapes.placeholders[54].text_frame
        page2_21_frame.text = data2_21
        page2_22_frame = shapes.placeholders[67].text_frame
        page2_22_frame.text = data2_22
        page2_23_frame = shapes.placeholders[68].text_frame
        page2_23_frame.text = data2_23
        page2_24_frame = shapes.placeholders[69].text_frame
        page2_24_frame.text = data2_24
        page2_25_frame = shapes.placeholders[70].text_frame
        page2_25_frame.text = data2_25
        page2_26_frame = shapes.placeholders[71].text_frame
        page2_26_frame.text = data2_26
        page2_27_frame = shapes.placeholders[47].text_frame
        page2_27_frame.text = data2_27
        page2_28_frame = shapes.placeholders[48].text_frame
        page2_28_frame.text = data2_28
        page2_29_frame = shapes.placeholders[27].text_frame
        page2_29_frame.text = data2_29
        page2_30_frame = shapes.placeholders[28].text_frame
        page2_30_frame.text = data2_30
        page2_31_frame = shapes.placeholders[66].text_frame
        page2_31_frame.text = data2_31
        page2_32_frame = shapes.placeholders[77].text_frame
        page2_32_frame.text = data2_32
        page2_33_frame = shapes.placeholders[78].text_frame
        page2_33_frame.text = data2_33
        page2_34_frame = shapes.placeholders[79].text_frame
        page2_34_frame.text = data2_34
        page2_35_frame = shapes.placeholders[80].text_frame
        page2_35_frame.text = data2_35
        page2_36_frame = shapes.placeholders[81].text_frame
        page2_36_frame.text = data2_36
        page2_37_frame = shapes.placeholders[82].text_frame
        page2_37_frame.text = data2_37
        page2_38_frame = shapes.placeholders[83].text_frame
        page2_38_frame.text = data2_38
        page2_39_frame = shapes.placeholders[84].text_frame
        page2_39_frame.text = data2_39
        page2_40_frame = shapes.placeholders[85].text_frame
        page2_40_frame.text = data2_40


        ##### 서비스 중단 여부 동그라미 #####
        circle_img_path = './images/circle.png'
        if script2_30 == '무중단 작업 가능':
            left_cm = 7.86
        elif script2_30 == '중단':
            left_cm = 9.92
        left = Inches(left_cm/2.54)   # inch = cm/2.54
        top_cm = 13.29
        top = Inches(top_cm/2.54)
        pic = slide.shapes.add_picture(circle_img_path, left, top)

        fig = (2, 1)
        plt.figure(figsize=fig)
        plt.barh(1, contract_speed_M, color='#00C0AA')
        plt.axvline(x=contract_speed_M, color='red', linewidth=0.5, linestyle='--')
        plt.barh(1, offer_speed_M, color='#4C4C4E', alpha=0.8)
        plt.axvline(x=offer_speed_M, color='red', linewidth=0.5, linestyle='--')
        plt.axis('off')
        plt.ylim(0, 2)
        image_stream_types = io.BytesIO()
        plt.savefig(image_stream_types, bbox_inches='tight', transparent=True)
        left = Inches(0.8)
        top = Inches(5.2)
        width = Inches(1.6)
        height = Inches(0.8)
        pic = shapes.add_picture(image_stream_types, left, top, width, height)
        plt.close()

        #####동종업계 속도 이용현황#####

        if bidw_df['contract_speed'][0][-1] == 'G':
            contract_sp = int(bidw_df['contract_speed'][0][:-1]) * pow(10, 9)
        elif bidw_df['contract_speed'][0][-1] == 'M':
            contract_sp = int(bidw_df['contract_speed'][0][:-1]) * pow(10, 6)

        if contract_sp >= 1000000000:
            explode = [0.1, 0, 0, 0]
        elif 1000000000 > contract_sp >= 500000000:
            explode = [0, 0.1, 0, 0]
        elif 500000000 > contract_sp >= 100000000:
            explode = [0, 0, 0.1, 0]
        elif 100000000 > contract_sp:
            explode = [0, 0, 0, 0.1]

        fig = (2.8, 3.3)
        plt.figure(figsize=fig)
        ratio = [over_1G, over_500M, over_100M, under_100M]
        plt.pie(ratio, autopct='%.1f%%', colors=['red', '#00C0AA', '#D1D2D4', '#4C4C4E'], counterclock=False,
                explode=explode)

        image_stream_types = io.BytesIO()
        plt.savefig(image_stream_types, bbox_inches='tight', transparent=True)
        left = Inches(5.8)
        top = Inches(4.9)
        width = Inches(2)
        height = Inches(1.9)
        pic = shapes.add_picture(image_stream_types, left, top, width, height)
        plt.close()

        #####동종업계 IP주소 자원 이용현황#####

        fig = (4, 4)
        plt.figure(figsize=fig)

        ip_list = [1000, 500, 256, 128, 64, 32, 16, 8, 4]
        colors = ['#e5e5e5', '#d4d4d4', '#c7c7c7', '#b6b6b6', '#979797', '#757575', '#595959', '#4c4c4c']
        for index, value in enumerate(ip_list):
            if total_ip >= value:
                colors.insert(index, '#00C0AA')
                break

        tick_list = ['1000', '500', '256', '128', '64', '32', '16', '8', '4']
        value = [ip_1000, ip_500, ip_256, ip_128, ip_64, ip_32,
                 ip_16, ip_8, ip_4]
        plt.bar(tick_list, value, color=colors)

        image_stream_types = io.BytesIO()
        plt.savefig(image_stream_types, bbox_inches='tight', transparent=True)
        left = Inches(8.1)
        top = Inches(5)
        width = Inches(2.3)
        height = Inches(2.1)
        pic = shapes.add_picture(image_stream_types, left, top, width, height)
        plt.close()
        print('세 번째 슬라이드 생성 완료')
        
        # img_path = './images/img.png'
        #
        # left = Inches(5.55)
        # top = Inches(1.5)
        # pic = slide.shapes.add_picture(img_path, left, top)

        # 추가 페이지

        if int(max(traffic_df['traffic'])) >= int(offer_speed*0.7):
            slide_layout = prs.slide_layouts[4]
            slide = prs.slides.add_slide(slide_layout)
            shapes = slide.shapes
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            shapes = slide.shapes
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            shapes = slide.shapes

        file_name = data1_1 + '_' + data1_3 + '_' + str(self.datetime) + '.pptx'
        url = './reports/' + file_name
        prs.save(url)

        return file_name
